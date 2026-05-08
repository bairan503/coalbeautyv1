import os
import json

# 尝试导入必要的库
try:
    from docx import Document
    docx_available = True
except ImportError:
    docx_available = False

try:
    from pptx import Presentation
    pptx_available = True
except ImportError:
    pptx_available = False

try:
    import PyPDF2
    pdf_available = True
except ImportError:
    pdf_available = False

def extract_docx(file_path):
    """提取docx文件内容"""
    if not docx_available:
        return "需要安装 python-docx 库: pip install python-docx"
    
    try:
        doc = Document(file_path)
        content = []
        for para in doc.paragraphs:
            if para.text.strip():
                content.append(para.text)
        return '\n'.join(content)
    except Exception as e:
        return f"读取失败: {str(e)}"

def extract_pptx(file_path):
    """提取pptx文件内容"""
    if not pptx_available:
        return "需要安装 python-pptx 库: pip install python-pptx"
    
    try:
        prs = Presentation(file_path)
        content = []
        for slide_num, slide in enumerate(prs.slides, 1):
            content.append(f"=== 幻灯片 {slide_num} ===")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
            content.append("")
        return '\n'.join(content)
    except Exception as e:
        return f"读取失败: {str(e)}"

def extract_pdf(file_path):
    """提取pdf文件内容"""
    if not pdf_available:
        return "需要安装 PyPDF2 库: pip install PyPDF2"
    
    try:
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            content = []
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text()
                if text and text.strip():
                    content.append(f"=== 第 {page_num} 页 ===")
                    content.append(text)
            return '\n'.join(content)
    except Exception as e:
        return f"读取失败: {str(e)}"

def main():
    files = {
        '结项报告书.pdf': extract_pdf,
        '煤美与共定稿(3).docx': extract_docx,
        '煤美与共(1).pptx': extract_pptx
    }
    
    results = {}
    for filename, extract_func in files.items():
        filepath = f'f:\\TraeCN\\meidiao\\{filename}'
        print(f"正在处理: {filename}")
        if os.path.exists(filepath):
            content = extract_func(filepath)
            results[filename] = content
            # 保存单独的文件
            output_filename = f'{os.path.splitext(filename)[0]}.txt'
            with open(output_filename, 'w', encoding='utf-8') as f:
                f.write(content)
            print(f"内容已保存到: {output_filename}")
        else:
            results[filename] = f"文件不存在: {filepath}"
            print(f"文件不存在: {filepath}")
    
    # 保存汇总结果
    with open('content_summary.json', 'w', encoding='utf-8') as f:
        json.dump(results, f, ensure_ascii=False, indent=2)
    print("\n汇总结果已保存到: content_summary.json")

if __name__ == '__main__':
    main()